#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""对生成的 pptx 做静态自查:
  1) 形状是否越出页面边界
  2) 估算每个文本框换行后的高度, 找出可能溢出其容器的文字
估算规则: 中日韩字符按 1.0 em 宽, 其余按 0.5 em; 行高 = size * line_spacing。
"""
import sys, math, unicodedata
from pptx import Presentation
from pptx.util import Emu, Pt

EMU_IN = 914400.0


def is_wide(ch):
    return unicodedata.east_asian_width(ch) in ("W", "F")


def text_width_em(s):
    return sum(1.0 if is_wide(c) else 0.5 for c in s)


def main(path):
    prs = Presentation(path)
    SW, SH = prs.slide_width, prs.slide_height
    problems = []
    for idx, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            # 1) 边界检查 (名字以 bleed 开头的是有意出血的装饰, 跳过)
            if sh.left is None or (sh.name or "").startswith("bleed"):
                continue
            l, t = sh.left, sh.top
            r, b = l + (sh.width or 0), t + (sh.height or 0)
            if l < -10000 or t < -10000 or r > SW + 10000 or b > SH + 10000:
                problems.append("p%02d 越界: %-22s L%.2f T%.2f R%.2f B%.2f in"
                                % (idx, sh.shape_type, l / EMU_IN, t / EMU_IN,
                                   r / EMU_IN, b / EMU_IN))
            # 2) 文字高度估算
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            box_w_in = (sh.width or 0) / EMU_IN
            box_h_in = (sh.height or 0) / EMU_IN
            total_h_pt = 0.0
            for p in tf.paragraphs:
                txt = "".join(r.text for r in p.runs)
                if not txt:
                    continue
                size = None
                for r_ in p.runs:
                    if r_.font.size:
                        size = r_.font.size.pt
                        break
                size = size or 18.0
                ls = p.line_spacing or 1.0
                sb = p.space_before.pt if p.space_before else 0.0
                sa = p.space_after.pt if p.space_after else 0.0
                avail_pt = box_w_in * 72.0
                need_pt = text_width_em(txt) * size
                lines = max(1, int(math.ceil(need_pt / avail_pt))) if avail_pt > 0 else 1
                total_h_pt += sb + lines * size * (ls if ls > 1.4 else 1.22) + sa
            if total_h_pt > box_h_in * 72.0 + 2:
                head = "".join(r.text for r in tf.paragraphs[0].runs)[:34]
                problems.append("p%02d 文字可能溢出: 需 %.2fin / 框高 %.2fin  «%s»"
                                % (idx, total_h_pt / 72.0, box_h_in, head))
    if problems:
        print("\n".join(problems))
        print("\n共 %d 处待看" % len(problems))
    else:
        print("OK: 无越界, 无明显文字溢出")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "deck/ClaudePet.pptx")
