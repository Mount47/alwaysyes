#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 ClaudePet 项目介绍 PPT(16:9, 12 页, 开源推广向)。

用法:
    python3 -m venv /tmp/pptxenv && /tmp/pptxenv/bin/pip install python-pptx
    /tmp/pptxenv/bin/python deck/build_deck.py

产物: deck/ClaudePet.pptx
"""

import os
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 主题

BG      = RGBColor(0x0F, 0x11, 0x15)
PANEL   = RGBColor(0x1A, 0x1E, 0x26)
PANEL2  = RGBColor(0x23, 0x28, 0x33)
LINE    = RGBColor(0x2E, 0x34, 0x40)
TEXT    = RGBColor(0xF5, 0xF6, 0xF8)
SUB     = RGBColor(0xB6, 0xBD, 0xC8)
MUTED   = RGBColor(0x79, 0x83, 0x92)
ACCENT  = RGBColor(0xE8, 0x73, 0x4A)   # Claude 橙, 品牌主色
RED     = RGBColor(0xFF, 0x4D, 0x3D)
ORANGE  = RGBColor(0xFF, 0xA0, 0x2E)
GRAY    = RGBColor(0x8E, 0x94, 0x9E)
BLUE    = RGBColor(0x5A, 0xA9, 0xE6)
GREEN   = RGBColor(0x4E, 0xD0, 0x7A)

F_LATIN = "Helvetica Neue"
F_CN    = "PingFang SC"
F_MONO  = "Menlo"

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.72)              # 左右留白
CW = W - 2 * M                # 内容宽度

# ---------------------------------------------------------------- 底层helper


def _set_fonts(run, latin, ea):
    """python-pptx 只设 latin 字体; 中文要显式写 a:ea, 否则 PPT 用默认宋体。"""
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    el_latin = rPr.find(qn("a:latin"))
    # 顺序必须是 latin -> ea -> cs, 所以倒着插
    for tag in ("a:cs", "a:ea"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.Element(qn(tag))
            el_latin.addnext(el)
        el.set("typeface", ea)


class Box(object):
    """包一层 text_frame, 逐段 add() 不用管第一段是否已用过。"""

    def __init__(self, slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        shape = slide.shapes.add_textbox(x, y, w, h)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        self.tf = tf
        self.shape = shape
        self._used = False

    def add(self, text, size=14, color=SUB, bold=False, mono=False,
            align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=1.3):
        if self._used:
            p = self.tf.add_paragraph()
        else:
            p = self.tf.paragraphs[0]
            self._used = True
        p.alignment = align
        p.line_spacing = line
        p.space_before = Pt(space_before)
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        _set_fonts(run, F_MONO if mono else F_LATIN, F_MONO if mono else F_CN)
        return p


def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.0):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    sh.text_frame.text = ""
    return sh


def slide_new(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])       # blank
    bg = rect(s, 0, 0, W, H, fill=BG, shape=MSO_SHAPE.RECTANGLE)
    bg.line.fill.background()
    return s


def header(slide, kicker, title, sub=None):
    """页眉: 小标签 + 大标题 + 橙色短横 (+ 可选副标题)"""
    b = Box(slide, M, Inches(0.52), CW, Inches(0.3))
    b.add(kicker, size=11, color=ACCENT, bold=True)
    b2 = Box(slide, M, Inches(0.84), CW, Inches(0.6))
    b2.add(title, size=30, color=TEXT, bold=True, line=1.1)
    bar = rect(slide, M, Inches(1.55), Inches(0.62), Inches(0.055),
               fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    bar.line.fill.background()
    if sub:
        b3 = Box(slide, M, Inches(1.74), CW, Inches(0.34))
        b3.add(sub, size=13.5, color=MUTED, line=1.25)


def card(slide, x, y, w, h, title, lines, accent=ACCENT, title_size=16,
         body_size=12.5, fill=PANEL, pad=Inches(0.26), badge=None, gap=6):
    """一张信息卡: 标题 + 若干正文行(每行前带小圆点由调用方决定)"""
    rect(slide, x, y, w, h, fill=fill, line=LINE, radius=0.06)
    ty = y + pad
    if badge is not None:
        bg_ = rect(slide, x + pad, ty, Inches(0.34), Inches(0.34),
                   fill=accent, radius=0.3)
        bt = Box(slide, x + pad, ty + Inches(0.055), Inches(0.34), Inches(0.3))
        bt.add(badge, size=12, color=BG, bold=True, align=PP_ALIGN.CENTER)
        tx = x + pad + Inches(0.46)
        tw = w - pad * 2 - Inches(0.46)
    else:
        tx, tw = x + pad, w - pad * 2
    tb = Box(slide, tx, ty, tw, Inches(0.4))
    tb.add(title, size=title_size, color=accent, bold=True, line=1.15)
    # 正文可用高度只扣顶部 padding —— 底部 padding 当作余量, 免得卡片被迫做得很高
    body = Box(slide, x + pad, ty + Inches(0.44), w - pad * 2, h - pad - Inches(0.44))
    for i, ln in enumerate(lines):
        if isinstance(ln, tuple):
            txt, col = ln
        else:
            txt, col = ln, SUB
        body.add(txt, size=body_size, color=col, space_before=0 if i == 0 else gap, line=1.35)
    return body


def code_block(slide, x, y, w, h, lines, title=None):
    rect(slide, x, y, w, h, fill=RGBColor(0x08, 0x0A, 0x0D), line=LINE, radius=0.05)
    # 左侧橙色竖条
    bar = rect(slide, x, y + Inches(0.12), Inches(0.045), h - Inches(0.24),
               fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    bar.line.fill.background()
    b = Box(slide, x + Inches(0.34), y + Inches(0.24), w - Inches(0.6), h - Inches(0.4))
    if title:
        b.add(title, size=10.5, color=MUTED, bold=True, space_after=6)
    for i, ln in enumerate(lines):
        if isinstance(ln, tuple):
            txt, col = ln
        else:
            txt, col = ln, RGBColor(0xD8, 0xDE, 0xE6)
        b.add(txt, size=12, color=col, mono=True,
              space_before=0 if (i == 0 and not title) else 5, line=1.3)
    return b


def ay_mark(slide, x, y, size, font_size, dot=RED, fill=PANEL2, fg=ACCENT):
    """画一个 AY 品牌图标(圆角方块 + AY 字 + 右上角状态点)"""
    rect(slide, x, y, size, size, fill=fill, line=LINE, radius=0.22)
    b = Box(slide, x, y + Emu(int(size * 0.26)), size, Emu(int(size * 0.5)))
    b.add("AY", size=font_size, color=fg, bold=True, align=PP_ALIGN.CENTER)
    if dot is not None:
        d = Emu(int(size * 0.24))
        c = rect(slide, x + size - d + Emu(int(size * 0.06)),
                 y - Emu(int(size * 0.06)), d, d,
                 fill=dot, shape=MSO_SHAPE.OVAL)
        c.line.color.rgb = BG
        c.line.width = Pt(1.5)
    return


def arrow(slide, x, y, w, h, color=ACCENT):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.shadow.inherit = False
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.adjustments[0] = 0.5
    a.adjustments[1] = 0.42
    return a


def footer_strip(slide, text, y=Inches(6.42), color=SUB, fill=PANEL, accent=ACCENT):
    h = Inches(0.62)
    rect(slide, M, y, CW, h, fill=fill, line=LINE, radius=0.14)
    bar = rect(slide, M, y + Inches(0.1), Inches(0.045), h - Inches(0.2),
               fill=accent, shape=MSO_SHAPE.RECTANGLE)
    bar.line.fill.background()
    b = Box(slide, M + Inches(0.32), y, CW - Inches(0.6), h, anchor=MSO_ANCHOR.MIDDLE)
    b.add(text, size=13, color=color, line=1.25)


def page_num(slide, n, total=12):
    b = Box(slide, W - M - Inches(1.2), H - Inches(0.42), Inches(1.2), Inches(0.25))
    b.add("%02d / %d" % (n, total), size=9.5, color=RGBColor(0x4A, 0x52, 0x5E),
          align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- 各页

def s01_cover(prs):
    s = slide_new(prs)
    # 背景装饰: 右侧一块淡淡的橙色圆
    glow = rect(s, W - Inches(4.2), Inches(-1.2), Inches(5.6), Inches(5.6),
                fill=RGBColor(0x1C, 0x18, 0x18), shape=MSO_SHAPE.OVAL)
    glow.line.fill.background()
    glow.name = "bleed-glow"          # 出血装饰, 自查脚本据此跳过边界检查

    ay_mark(s, M, Inches(1.62), Inches(1.15), 34)

    b = Box(s, M, Inches(3.05), Inches(9.5), Inches(1.1))
    b.add("ClaudePet", size=58, color=TEXT, bold=True, line=1.0)

    b2 = Box(s, M, Inches(4.12), Inches(9.6), Inches(0.5))
    b2.add("给 Claude Code 的状态灯 + 桌面宠物", size=22, color=ACCENT, bold=True)

    b3 = Box(s, M, Inches(4.74), Inches(9.2), Inches(0.5))
    b3.add("AI 停下来等你点 yes 的那一刻 —— 让你一眼看见,再点一下跳回去。",
           size=15, color=MUTED, line=1.4)

    chips = ["macOS 原生", "Swift / AppKit", "无 Electron · 几 MB", "hook 驱动 · 不联网"]
    cx = M
    for c in chips:
        w = Inches(0.34 + 0.115 * len(c))
        rect(s, cx, Inches(5.72), w, Inches(0.42), fill=PANEL, line=LINE, radius=0.4)
        cb = Box(s, cx, Inches(5.72), w, Inches(0.42), anchor=MSO_ANCHOR.MIDDLE)
        cb.add(c, size=11.5, color=SUB, align=PP_ALIGN.CENTER)
        cx += w + Inches(0.16)

    b4 = Box(s, M, Inches(6.62), CW, Inches(0.3))
    b4.add("一个很小的痛点,一条很短的状态链路,和一堆把它做到「不烦人」的细节。",
           size=12, color=RGBColor(0x5A, 0x62, 0x6E))
    notes(s, "开场一句话: 这不是一个玩具桌宠, 它解决的是多项目并行跑 Claude Code 时"
             "「AI 在等你、你却不知道」的空转问题。桌宠只是外壳。")
    return s


def s02_problem(prs):
    s = slide_new(prs)
    header(s, "PROBLEM", "问题不是 AI 太慢,是你不知道它在等你")
    y = Inches(2.15)
    h = Inches(3.0)
    w = (CW - Inches(0.44) * 2) / 3
    items = [
        ("1", "你同时开着 3 个项目",
         ["Claude Code 在跑长任务,你顺手切去看文档、开会、刷手机。",
          "终端沉到后台,你完全不看它。"]),
        ("2", "它停在一句 Allow? 上",
         ["工具要授权、命令要确认 —— 它停下来等你。",
          "没有声音,没有角标,没有任何东西告诉你。"]),
        ("3", "20 分钟后你才发现",
         ["回头一看,进度条早就卡在那一步。",
          "并行的项目越多,这种空转就越久。"]),
    ]
    for i, (n, t, lines) in enumerate(items):
        x = M + i * (w + Inches(0.44))
        card(s, x, y, w, h, t, lines, badge=n, title_size=16.5, body_size=13)
        if i < 2:
            arrow(s, x + w + Inches(0.09), y + Inches(1.3), Inches(0.26), Inches(0.3),
                  color=RGBColor(0x3A, 0x41, 0x4D))
    footer_strip(s, "「等 AI」于是变成了「AI 等你」。缺的不是模型能力,是一盏状态灯。")
    page_num(s, 2)
    notes(s, "这一页要让听众自己点头。重点是「并行度越高越痛」——"
             "单开一个终端你盯着就行, 三个以上必然漏。")
    return s


def s03_what(prs):
    s = slide_new(prs)
    header(s, "WHAT IT IS", "两个东西:菜单栏一盏灯,桌面一只宠物")
    y = Inches(2.15)
    h = Inches(3.95)
    w = (CW - Inches(0.5)) / 2

    rect(s, M, y, w, h, fill=PANEL, line=LINE, radius=0.05)
    ay_mark(s, M + Inches(0.34), y + Inches(0.32), Inches(0.7), 20)
    tb = Box(s, M + Inches(1.2), y + Inches(0.42), w - Inches(1.5), Inches(0.4))
    tb.add("菜单栏 AY 图标", size=17, color=ACCENT, bold=True)
    bb = Box(s, M + Inches(0.34), y + Inches(1.3), w - Inches(0.68), h - Inches(1.6))
    for i, t in enumerate([
        "常驻菜单栏,不占地方。用颜色 + 红点编码所有会话的汇总状态。",
        "有项目等你确认时变红并带数字,同时弹一条 macOS 桌面通知。",
        "下拉逐条列出「项目名 — 状态 (等了多久)」,waiting 的排最前。",
        "每一行都能点 —— 点了就把那个会话所在的终端切到最前。",
    ]):
        bb.add("· " + t, size=13, color=SUB, space_before=0 if i == 0 else 11, line=1.4)

    x2 = M + w + Inches(0.5)
    rect(s, x2, y, w, h, fill=PANEL, line=LINE, radius=0.05)
    eb = rect(s, x2 + Inches(0.34), y + Inches(0.32), Inches(0.7), Inches(0.7),
              fill=PANEL2, line=LINE, radius=0.22)
    et = Box(s, x2 + Inches(0.34), y + Inches(0.44), Inches(0.7), Inches(0.5))
    et.add("🐣", size=24, color=TEXT, align=PP_ALIGN.CENTER)
    tb2 = Box(s, x2 + Inches(1.2), y + Inches(0.42), w - Inches(1.5), Inches(0.4))
    tb2.add("桌面浮动宠物", size=17, color=ACCENT, bold=True)
    bb2 = Box(s, x2 + Inches(0.34), y + Inches(1.3), w - Inches(0.68), h - Inches(1.6))
    for i, t in enumerate([
        "透明背景的小窗口,可以拖到任意屏幕、任意全屏空间,位置记住。",
        "跟着会话状态换逐帧动画:等确认 / 运行中 / 空闲 / 出错各有形象。",
        "形象可以换 —— 从社区画廊 petdex 装一只精灵图宠物。",
        "不想要就在菜单里「隐藏桌面宠物」,只留菜单栏那盏灯。",
    ]):
        bb2.add("· " + t, size=13, color=SUB, space_before=0 if i == 0 else 11, line=1.4)

    footer_strip(s, "设计原则:低互动、被动指示 —— 它不抢你的注意力,只在该抢的那一刻抢。",
                 y=Inches(6.42))
    page_num(s, 3)
    notes(s, "强调两者是一套数据的两个出口: 同一份状态, 一个给「不想被打扰」的人, "
             "一个给「想要点乐子」的人。")
    return s


def s04_states(prs):
    s = slide_new(prs)
    header(s, "STATES", "一眼读懂:颜色和动画就是全部语义")
    rows = [
        ("菜单栏 AY", "宠物动画", "含义", None, True),
        ("AY 变红 + 红点(多个跟数字)", "review  🐔", "有项目在等你点 yes", RED, False),
        ("AY 变橙 + 橙点", "failed  💥", "有任务出错", ORANGE, False),
        ("AY + 灰点", "run  🐥", "有会话正在跑", GRAY, False),
        ("AY 模板色(随深浅色菜单栏)", "idle  🐣", "全部空闲 / 已停用", SUB, False),
    ]
    top = Inches(2.2)
    rh = Inches(0.58)
    tw = CW
    cols = [Inches(4.5), Inches(3.0), tw - Inches(7.5)]
    for r, (c1, c2, c3, col, is_head) in enumerate(rows):
        y = top + r * rh
        fill = PANEL2 if is_head else (PANEL if r % 2 else RGBColor(0x15, 0x18, 0x1F))
        bar = rect(s, M, y, tw, rh - Inches(0.04), fill=fill,
                   shape=MSO_SHAPE.RECTANGLE)
        bar.line.fill.background()
        if not is_head:
            dot = rect(s, M + Inches(0.26), y + Inches(0.21), Inches(0.16), Inches(0.16),
                       fill=col, shape=MSO_SHAPE.OVAL)
            dot.line.fill.background()
        x = M + Inches(0.58)
        for ci, txt in enumerate((c1, c2, c3)):
            cb = Box(s, x, y, cols[ci], rh - Inches(0.04), anchor=MSO_ANCHOR.MIDDLE)
            cb.add(txt, size=11 if is_head else 13.5,
                   color=MUTED if is_head else (col if ci == 0 else SUB),
                   bold=is_head or ci == 0)
            x += cols[ci]

    y2 = top + len(rows) * rh + Inches(0.22)
    w = (CW - Inches(0.4)) / 2
    card(s, M, y2, w, Inches(1.5), "状态跃迁会插播一次性动画",
         ["新冒出等你确认的项目 → jump 跳一下引起注意",
          "手头的活全干完 → wave 挥手庆祝 2 秒再回常态"], accent=GREEN, body_size=12.5)
    card(s, M + w + Inches(0.4), y2, w, Inches(1.5), "failed 已就绪,但暂时没人写它",
         ["菜单栏橙色 + failed 动画,app 侧都已做好",
          "缺的是一个 PostToolUse hook 去判错 —— 挂上即生效"],
         accent=ORANGE, body_size=12.5)
    page_num(s, 4)
    notes(s, "对照表是全片信息密度最高的一页, 讲的时候只念前两行, 剩下让人扫一眼即可。")
    return s


def s05_install(prs):
    s = slide_new(prs)
    header(s, "INSTALL", "clone,一条命令,它就在菜单栏了")
    y = Inches(2.15)
    w1 = Inches(7.0)
    code_block(s, M, y, w1, Inches(2.2), [
        ("git clone <repo> && cd alwaysyes", TEXT),
        ("", TEXT),
        ("./install.sh", GREEN),
        ("    # 装 hook 脚本 + 合并配置 + 构建 app", MUTED),
        ("bin/claudepet start", GREEN),
        ("    # 启动,菜单栏出现 AY 图标", MUTED),
    ], title="TERMINAL")

    card(s, M, y + Inches(2.42), w1, Inches(1.83), "install.sh 幂等地做四件事",
         ["① 复制 hooks/ 下的脚本到 ~/.claude/hooks/",
          "② 用 jq 合并 hooks 段进 settings.json(自动备份)",
          "③ 生成默认 ~/.claude/pet-config.json(已有则保留)",
          "④ 用 swiftc 构建 ClaudePet.app"], body_size=11.5, accent=BLUE, gap=5)

    x2 = M + w1 + Inches(0.44)
    w2 = CW - w1 - Inches(0.44)
    body = card(s, x2, y, w2, Inches(4.25), "claudepet — 命令行入口", [], accent=ACCENT)
    cmds = [
        ("claudepet start", "启动桌面宠物"),
        ("claudepet stop", "退出"),
        ("claudepet status", "命令行看各项目状态"),
        ("claudepet install", "安装 hooks 并首次构建"),
        ("claudepet pet --list", "列出可装的宠物形象"),
        ("claudepet pet <slug>", "装一只宠物"),
    ]
    cy = y + Inches(0.86)
    for c, d in cmds:
        cb = Box(s, x2 + Inches(0.3), cy, w2 - Inches(0.6), Inches(0.28))
        cb.add(c, size=12, color=GREEN, mono=True)
        cb2 = Box(s, x2 + Inches(0.3), cy + Inches(0.26), w2 - Inches(0.6), Inches(0.26))
        cb2.add(d, size=11, color=MUTED)
        cy += Inches(0.56)

    footer_strip(s, "依赖只有 jq / swiftc / curl。没装 swiftc 也不会失败 —— 只是跳过构建,"
                    "其余照装。hook 对新会话生效,已开着的 claude 需要重启。")
    page_num(s, 5)
    notes(s, "演示重点: 三行命令就完事, 而且 install.sh 不会破坏你已有的 settings.json。")
    return s


def s06_pet(prs):
    s = slide_new(prs)
    header(s, "INTERACTION", "它在桌面上,但绝不挡事")
    y = Inches(2.15)
    w = (CW - Inches(0.44)) / 2
    h = Inches(1.95)
    card(s, M, y, w, h, "点击穿透", [
        "鼠标只在画到实际像素的地方被拦截。",
        "宠物周围的透明区域点下去,事件会落到底下的窗口 —— 它压在代码上也不碍事。",
    ], accent=BLUE, body_size=13)
    card(s, M + w + Inches(0.44), y, w, h, "拖到哪都行", [
        "按住就拖,多显示器、全屏空间随便摆。",
        "位置会记住,下次启动还在原地。",
    ], accent=GREEN, body_size=13)
    y2 = y + h + Inches(0.34)
    card(s, M, y2, w, h, "气泡按需弹", [
        "默认只用动画表达状态,不常驻文字。",
        "鼠标悬停或左键点一下才弹气泡,逐行列出全部正在等待的项目,移开即收。",
    ], accent=ACCENT, body_size=13)
    card(s, M + w + Inches(0.44), y2, w, h, "右键就是控制台", [
        "隐藏宠物 / 刷新状态 (R) / 重置状态 / 退出。",
        "菜单栏那份下拉菜单里也有同样的开关。",
    ], accent=RGBColor(0xB4, 0x8E, 0xFF), body_size=13)
    footer_strip(s, "宠物的每一处交互都在回答同一个问题:怎么让它存在感刚好够,又不成为负担。")
    page_num(s, 6)
    notes(s, "点击穿透是最容易被低估的一条 —— 没有它, 一只永远置顶的宠物就是障碍物。")
    return s


def s07_jump(prs):
    s = slide_new(prs)
    header(s, "JUMP BACK", "点一下项目名,直接切回那个终端",
           "hook 会把会话的 tty、$TERM_PROGRAM 和 cwd 一起写进状态文件,app 按可用信息分三档定位。")
    tiers = [
        ("精确到标签页", RED,
         "iTerm2 · Apple Terminal",
         "它们的 AppleScript 词典里标签页带 tty 属性 —— 按 tty 直接选中那一个 tab。"),
        ("提到最前", ORANGE,
         "Ghostty · WezTerm · kitty · Alacritty · Warp",
         "没有可脚本化的标签页模型,退化为把这个 app 整体拉到最前。"),
        ("反推窗口", BLUE,
         "VS Code · Cursor 插件里的会话",
         "跑在 stdio 上,既没 tty 也没 TERM_PROGRAM。靠 cwd 命中 ~/.claude/ide/*.lock 里的 "
         "workspaceFolders,反推出是哪个 IDE 再激活它。"),
    ]
    y = Inches(2.4)
    h = Inches(1.16)
    for i, (label, col, apps, desc) in enumerate(tiers):
        yy = y + i * (h + Inches(0.22))
        rect(s, M, yy, CW, h, fill=PANEL, line=LINE, radius=0.08)
        bar = rect(s, M, yy + Inches(0.12), Inches(0.05), h - Inches(0.24),
                   fill=col, shape=MSO_SHAPE.RECTANGLE)
        bar.line.fill.background()
        lb = Box(s, M + Inches(0.34), yy + Inches(0.2), Inches(2.0), Inches(0.4))
        lb.add(label, size=15.5, color=col, bold=True)
        lb2 = Box(s, M + Inches(0.34), yy + Inches(0.66), Inches(2.4), Inches(0.3))
        lb2.add("第 %d 档" % (i + 1), size=10.5, color=MUTED)
        ab = Box(s, M + Inches(2.7), yy + Inches(0.2), Inches(5.4), Inches(0.34))
        ab.add(apps, size=12.5, color=TEXT, bold=True, mono=False)
        db = Box(s, M + Inches(2.7), yy + Inches(0.6), CW - Inches(3.1), Inches(0.5))
        db.add(desc, size=12, color=SUB, line=1.35)
    footer_strip(s, "首次点击时 macOS 会弹一次「允许 ClaudePet 控制 iTerm」的自动化授权,同意一次即可。"
                    "(ad-hoc 签名,重新构建后可能再问一次)")
    page_num(s, 7)
    notes(s, "这页想传达的是「诚实」: 不吹能 100% 命中, 而是说清哪种终端能到什么精度。")
    return s


def s08_skins(prs):
    s = slide_new(prs)
    header(s, "SKINS", "形象来自社区画廊 petdex")
    y = Inches(2.15)
    w1 = Inches(5.9)
    code_block(s, M, y, w1, Inches(1.9), [
        ("claudepet pet --list cat", GREEN),
        ("    # 按关键词浏览可用形象", MUTED),
        ("claudepet pet homelander", GREEN),
        ("    # 装到 ~/.claude/pets/<slug>/", MUTED),
    ], title="装一只")
    card(s, M, y + Inches(2.1), w1, Inches(2.15), "官方渠道装的也能直接选", [
        "app 会扫三个目录,不只是自己装的那个:",
        ("~/.claude/pets/          本脚本", GREEN),
        ("~/.codex/pets/           Codex / awesome-codex-pet", GREEN),
        ("~/.petdex/pets/          petdex CLI", GREEN),
        "所以 npx petdex install boba 装的宠物,在菜单里一样能选中。",
    ], accent=BLUE, body_size=12)

    x2 = M + w1 + Inches(0.44)
    w2 = CW - w1 - Inches(0.44)
    card(s, x2, y, w2, Inches(4.25), "精灵图规范:v1 / v2 都吃", [
        "v1 = 8 列 × 9 行 (1536×1872)",
        "v2 = 8 列 × 11 行 (1536×2288),多出的两行是 16 个朝向,本 app 暂不用",
        "单帧恒为 192×208px;行 = 动画状态 (idle / wave / run / failed / review / jump),每态 6 帧循环",
        ("行数按帧宽高比自动反推,不写死 —— 以后出 v3 也不用改代码。", TEXT),
        ("找不到精灵图时自动回退到 emoji 占位,不会白屏。", TEXT),
    ], accent=ACCENT, body_size=12.5)

    footer_strip(s, "版权:petdex 的美术资源版权归各自提交者所有(含同人作品)。"
                    "ClaudePet 只通过其公开 API 下载到用户本地,不打包、不转发任何精灵图。",
                 color=MUTED, accent=RGBColor(0x5A, 0x62, 0x6E))
    page_num(s, 8)
    notes(s, "复用现成生态而不是自造格式 —— 这也是为什么装机第一天就有几十只宠物可选。")
    return s


def s09_how(prs):
    s = slide_new(prs)
    header(s, "HOW IT WORKS", "hook 写状态,app 读状态 —— 就这么点东西")
    y = Inches(2.22)
    h = Inches(3.35)
    w1 = Inches(4.25)
    w3 = Inches(4.25)
    gap = Inches(0.42)
    w2 = CW - w1 - w3 - gap * 2

    hooks = [
        ("SessionStart", "→ idle", GREEN),
        ("UserPromptSubmit", "→ running", GRAY),
        ("Notification", "→ waiting + 通知", RED),
        ("Stop", "→ idle", GREEN),
        ("SessionEnd", "→ 删除状态文件", MUTED),
    ]
    rect(s, M, y, w1, h, fill=PANEL, line=LINE, radius=0.06)
    hb = Box(s, M + Inches(0.28), y + Inches(0.26), w1 - Inches(0.56), Inches(0.4))
    hb.add("Claude Code hooks", size=15, color=ACCENT, bold=True)
    hb2 = Box(s, M + Inches(0.28), y + Inches(0.7), w1 - Inches(0.56), Inches(0.3))
    hb2.add("5 个几十行的 shell 脚本", size=11, color=MUTED)
    hy = y + Inches(1.16)
    for name, act, col in hooks:
        r = Box(s, M + Inches(0.28), hy, Inches(2.0), Inches(0.28))
        r.add(name, size=11.5, color=TEXT, mono=True)
        r2 = Box(s, M + Inches(2.28), hy, w1 - Inches(2.56), Inches(0.28))
        r2.add(act, size=11.5, color=col)
        hy += Inches(0.42)

    x2 = M + w1 + gap
    arrow(s, x2 - Inches(0.36), y + Inches(1.5), Inches(0.3), Inches(0.34))
    rect(s, x2, y, w2, h, fill=RGBColor(0x08, 0x0A, 0x0D), line=ACCENT, radius=0.06, line_w=1.2)
    sb = Box(s, x2 + Inches(0.26), y + Inches(0.26), w2 - Inches(0.52), Inches(0.4))
    sb.add("状态文件", size=15, color=ACCENT, bold=True)
    sb2 = Box(s, x2 + Inches(0.26), y + Inches(0.72), w2 - Inches(0.52), Inches(0.7))
    sb2.add("~/.claude/pet-state/\n  <session_id>.json", size=11.5,
            color=GREEN, mono=True, line=1.35)
    sb3 = Box(s, x2 + Inches(0.26), y + Inches(1.62), w2 - Inches(0.52), Inches(1.72))
    for i, t in enumerate([
        "一个会话一个文件,天然避开并发写冲突。",
        "字段:project / cwd / status / tty / term_program / pid / updated_at",
        "纯本地纯文本,你随时能 cat 出来看。",
    ]):
        sb3.add("· " + t, size=11.5, color=SUB, space_before=0 if i == 0 else 9, line=1.35)

    x3 = x2 + w2 + gap
    arrow(s, x3 - Inches(0.36), y + Inches(1.5), Inches(0.3), Inches(0.34))
    rect(s, x3, y, w3, h, fill=PANEL, line=LINE, radius=0.06)
    ab = Box(s, x3 + Inches(0.28), y + Inches(0.26), w3 - Inches(0.56), Inches(0.4))
    ab.add("ClaudePet.app", size=15, color=ACCENT, bold=True)
    ab2 = Box(s, x3 + Inches(0.28), y + Inches(0.7), w3 - Inches(0.56), Inches(0.3))
    ab2.add("~2100 行 Swift / AppKit", size=11, color=MUTED)
    ab3 = Box(s, x3 + Inches(0.28), y + Inches(1.16), w3 - Inches(0.56), Inches(2.0))
    for i, t in enumerate([
        "FSEvents 实时监听该目录,5 秒兜底刷新。",
        "扫进程对账补漏(ps + lsof + argv 里的 --resume)。",
        "汇总所有会话 → 菜单栏图标 + 桌面宠物 + 通知。",
    ]):
        ab3.add("· " + t, size=11.5, color=SUB, space_before=0 if i == 0 else 10, line=1.35)

    footer_strip(s, "没有守护进程,没有服务端,不联网,也不读你的代码内容 —— 只读它自己写的那几个 JSON。")
    page_num(s, 9)
    notes(s, "架构简单是卖点不是缺点: 整条链路你 5 分钟能读完, 出问题自己就能查。")
    return s


def s10_safe(prs):
    s = slide_new(prs)
    header(s, "SAFE BY DESIGN", "装它这件事本身,不该让你有风险")
    y = Inches(2.15)
    w = (CW - Inches(0.44)) / 2
    h = Inches(1.98)
    card(s, M, y, w, h, "不覆盖你的配置", [
        "install.sh 幂等,跑几次都一样。",
        "用 jq 把 hooks 段合并进 settings.json:先滤掉指向本项目的旧条目再追加,"
        "你原有的 hook / env / permissions 全部原样保留,并自动备份。",
    ], accent=GREEN, body_size=12.5)
    card(s, M + w + Inches(0.44), y, w, h, "仓库 = 唯一事实来源", [
        "git 里只有 hooks/*.sh 与 Swift 源码。",
        "~/.claude/ 下的全是安装产物,不入库 —— 含 token 的 settings.json 永远不会被提交。",
    ], accent=BLUE, body_size=12.5)
    y2 = y + h + Inches(0.34)
    card(s, M, y2, w, h, "三层开关,逐层更硬", [
        "① 菜单项:停用(全部) / 隐藏桌面宠物 —— 最方便",
        "② ~/.claude/pet-config.json:单一事实来源",
        "③ settings.json 里摘掉 hook —— 最彻底,连状态都不再写",
    ], accent=ACCENT, body_size=12.5)
    card(s, M + w + Inches(0.44), y2, w, h, "状态会自愈", [
        "进程判活优先(kill -0),进程没了立刻清;SessionEnd 正常结束直接删。",
        "没记 pid 的老文件走超时兜底。菜单里「刷新状态」还能扫进程对账 —— "
        "但 hook 写的状态永远优先,扫描只补缺和清死。",
    ], accent=RGBColor(0xB4, 0x8E, 0xFF), body_size=12.5)
    footer_strip(s, "已经用假的 settings.json 模拟过「别人的电脑」场景,确认不会破坏既有配置。")
    page_num(s, 10)
    notes(s, "开源推广最大的阻力是「我不敢让一个脚本去改我的 settings.json」。这页就是拆这颗雷。")
    return s


def s11_footprint(prs):
    s = slide_new(prs)
    header(s, "FOOTPRINT", "几 MB,一个进程,零网络请求")
    y = Inches(2.2)
    stats = [
        ("~2100", "行源码,全部读完不用半小时"),
        ("6", "个 Swift 文件 + 5 个 hook 脚本"),
        ("3", "个依赖:jq / swiftc / curl"),
        ("0", "个后台服务、0 次联网上报"),
    ]
    w = (CW - Inches(0.36) * 3) / 4
    for i, (n, d) in enumerate(stats):
        x = M + i * (w + Inches(0.36))
        rect(s, x, y, w, Inches(1.72), fill=PANEL, line=LINE, radius=0.08)
        nb = Box(s, x, y + Inches(0.3), w, Inches(0.75))
        nb.add(n, size=42, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, line=1.0)
        db = Box(s, x + Inches(0.18), y + Inches(1.08), w - Inches(0.36), Inches(0.55))
        db.add(d, size=11.5, color=SUB, align=PP_ALIGN.CENTER, line=1.3)

    y2 = y + Inches(2.05)
    w2 = (CW - Inches(0.44)) / 2
    card(s, M, y2, w2, Inches(1.85), "为什么是 Swift / AppKit,不是 Electron", [
        "不打包 Chromium,进程常驻内存只有几 MB。",
        "LSUIElement = true,纯菜单栏应用,不占 Dock、不抢焦点。",
        "系统原生渲染,菜单栏图标能跟着深浅色主题自动适配。",
    ], accent=GREEN, body_size=12.5)
    card(s, M + w2 + Inches(0.44), y2, w2, Inches(1.85), "代价:仅 macOS", [
        "重度依赖 Cocoa / AppKit 和 AppleScript,没法直接跨平台。",
        "需要 Windows / Linux,或者想联动实体硬件桌宠,可以看看",
        ("Seeed-Solution/vibe-pet", BLUE),
    ], accent=ORANGE, body_size=12.5)
    page_num(s, 11)
    notes(s, "把「只支持 mac」当作明确的取舍讲出来, 比藏着掖着更容易赢得信任。")
    return s


def s12_next(prs):
    s = slide_new(prs)
    header(s, "ROADMAP", "接下来想做的")
    y = Inches(2.15)
    w = (CW - Inches(0.44)) / 2
    items_l = [
        "用量环 —— 读 statusline 的 rate_limits,在宠物外圈画 5h / 7d 额度",
        "气泡里直接 Allow / Deny(要接 permission hook)",
        "贴边 mini 模式:拖到屏幕边缘半隐藏,hover 探头",
        "完成 / 等待音效 + 免打扰开关",
    ]
    items_r = [
        "PostToolUse 判错自动变脸(failed 那套 app 侧已就绪)",
        "60 秒没鼠标动静就睡觉,鼠标一动惊醒",
        "区分「当前活跃会话」与「后台等待项目」",
        "claudepet refresh 命令行版 + Homebrew 分发",
    ]
    for ci, items in enumerate((items_l, items_r)):
        x = M + ci * (w + Inches(0.44))
        for i, t in enumerate(items):
            yy = y + i * Inches(0.7)
            rect(s, x, yy, w, Inches(0.58), fill=PANEL, line=LINE, radius=0.14)
            d = rect(s, x + Inches(0.24), yy + Inches(0.21), Inches(0.14), Inches(0.14),
                     fill=ACCENT, shape=MSO_SHAPE.OVAL)
            d.line.fill.background()
            b = Box(s, x + Inches(0.56), yy, w - Inches(0.8), Inches(0.58),
                    anchor=MSO_ANCHOR.MIDDLE)
            b.add(t, size=12, color=SUB, line=1.25)

    y2 = Inches(5.15)
    rect(s, M, y2, CW, Inches(1.55), fill=PANEL2, line=ACCENT, radius=0.1, line_w=1.2)
    cb = Box(s, M + Inches(0.4), y2 + Inches(0.26), Inches(6.0), Inches(0.45))
    cb.add("clone 下来就能跑", size=21, color=TEXT, bold=True)
    cb2 = Box(s, M + Inches(0.4), y2 + Inches(0.78), Inches(6.4), Inches(0.5))
    cb2.add("装完立刻见效,不喜欢就菜单里一键停用,或者直接摘掉 hook。",
            size=12.5, color=MUTED, line=1.3)
    code_block(s, M + Inches(6.9), y2 + Inches(0.2), CW - Inches(7.3), Inches(1.15), [
        ("git clone <repo> && cd alwaysyes", TEXT),
        ("./install.sh && bin/claudepet start", GREEN),
    ])
    page_num(s, 12)
    notes(s, "收尾把 CTA 放大: 门槛低 + 退出成本低, 是让人愿意试的两个关键。")
    return s


# ---------------------------------------------------------------- main

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for fn in (s01_cover, s02_problem, s03_what, s04_states, s05_install, s06_pet,
               s07_jump, s08_skins, s09_how, s10_safe, s11_footprint, s12_next):
        fn(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ClaudePet.pptx")
    prs.save(out)
    print("saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))


if __name__ == "__main__":
    main()
